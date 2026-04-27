"""
semantic_search.py — поиск по содержимому файлов через OpenAI Embeddings.

Стратегии извлечения текста:
  FULL_TEXT   — читаем файл напрямую (txt, md, код)
  LIBRARY     — извлекаем через библиотеку (docx, pdf)
  STRUCTURAL  — структурное описание без чтения данных (xlsx, csv, pptx)
  SKIP        — пропускаем (бинарники, конфиги, секреты)

Хранение: SQLite BLOB (float32 × 1536 = 6 КБ на файл)
Поиск:    numpy cosine similarity — достаточно до ~50k файлов без FAISS
"""

import csv
import os
import pathlib
import sqlite3
import struct
import threading
import time
import queue

import numpy as np

# ── Константы ─────────────────────────────────────────────────────────────────

DB_PATH     = pathlib.Path(__file__).parent / "semantic.db"
EMBED_MODEL = "text-embedding-3-small"
EMBED_DIM   = 1536
BATCH_SIZE  = 100       # файлов за один API-запрос
MAX_CHARS   = 3000      # символов текста на файл
MAX_FILE_MB = 500       # пропускаем только огромные файлы (>500 МБ)
SIM_THRESH  = 0.25      # минимальный cosine score

# ── Таблица расширений ────────────────────────────────────────────────────────

# Стратегия 1: читаем текст напрямую
FULL_TEXT: set[str] = {
    "txt", "md", "rst", "log",
    "py", "js", "ts", "tsx", "jsx", "vue",
    "java", "kt", "cpp", "c", "h", "cs", "go", "rs", "rb", "php", "swift",
    "html", "css", "scss", "sass", "sql", "sh", "bash", "ps1",
}

# Стратегия 2: извлекаем через библиотеку
LIBRARY: dict[str, str] = {
    "pdf":  "pymupdf",    # pip install pymupdf
    "docx": "python-docx",  # pip install python-docx
}

# Стратегия 3: структурное описание (имя + заголовки, без данных)
STRUCTURAL: dict[str, str] = {
    "xlsx": "openpyxl",    # pip install openpyxl
    "xls":  "openpyxl",
    "csv":  "builtin",
    "pptx": "python-pptx", # pip install python-pptx


}

# Никогда не индексируем
SKIP: set[str] = {
    # Конфиги и секреты — могут содержать пароли
    "env", "json", "xml", "yaml", "yml", "ini", "toml", "cfg", "conf",
    # Бинарники
    "exe", "dll", "so", "dylib", "bin", "dat",
    # Медиа
    "jpg", "jpeg", "png", "gif", "bmp", "webp", "svg",
    "mp3", "mp4", "avi", "mkv", "wav", "flac",
    # Архивы
    "zip", "rar", "7z", "tar", "gz",
    # Компилированный код
    "pyc", "class", "o", "obj",
}

ALL_SUPPORTED = FULL_TEXT | set(LIBRARY) | set(STRUCTURAL)


# ── Извлечение текста ─────────────────────────────────────────────────────────

def _too_large(path: str) -> bool:
    try:
        return os.path.getsize(path) > MAX_FILE_MB * 1024 * 1024
    except OSError:
        return True


def _pages_for_size(path: str) -> int:
    """Адаптивное число страниц в зависимости от размера файла."""
    try:
        mb = os.path.getsize(path) / (1024 * 1024)
    except OSError:
        return 1
    if mb < 5:   return 10
    if mb < 20:  return 3
    if mb < 100: return 2
    return 1


def _extract_text(path: str) -> str:
    """Возвращает текст файла по стратегии расширения. '' если не удалось."""
    ext = pathlib.Path(path).suffix.lower().lstrip(".")

    if ext in SKIP or ext not in ALL_SUPPORTED:
        return ""
    if _too_large(path):
        return ""

    # ── Стратегия 1: прямое чтение ────────────────────────────────────────────
    if ext in FULL_TEXT:
        for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
            try:
                with open(path, "r", encoding=enc, errors="strict") as f:
                    return f.read(MAX_CHARS)
            except (UnicodeDecodeError, LookupError):
                continue
            except Exception:
                return ""
        return ""

    # ── Стратегия 2: PDF ──────────────────────────────────────────────────────
    if ext == "pdf":
        max_pages = _pages_for_size(path)
        try:
            import pymupdf  # pymupdf >= 1.24
            doc = pymupdf.open(path)
            text = ""
            for page in doc[:max_pages]:
                text += page.get_text()
                if len(text) >= MAX_CHARS:
                    break
            doc.close()
            return text[:MAX_CHARS]
        except ImportError:
            pass
        except Exception:
            return ""

        # Fallback: pypdf
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            text = ""
            for page in reader.pages[:max_pages]:
                text += page.extract_text() or ""
                if len(text) >= MAX_CHARS:
                    break
            return text[:MAX_CHARS]
        except ImportError:
            pass
        except Exception:
            pass
        return ""

    # ── Стратегия 2: DOCX ─────────────────────────────────────────────────────
    if ext == "docx":
        try:
            import docx as _docx
            doc = _docx.Document(path)
            text = " ".join(
                p.text for p in doc.paragraphs if p.text.strip()
            )
            return text[:MAX_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    # ── Стратегия 3: XLSX ─────────────────────────────────────────────────────
    if ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            filename = pathlib.Path(path).name
            lines = [f"Excel файл: {filename}"]
            for sheet_name in wb.sheetnames[:10]:
                ws = wb[sheet_name]
                # Берём только первую строку (заголовки) — данные не нужны
                first_row = next(ws.iter_rows(max_row=1, values_only=True), ())
                headers = [str(v) for v in first_row if v is not None]
                if headers:
                    lines.append(f"Лист '{sheet_name}': {', '.join(headers)}")
                else:
                    lines.append(f"Лист '{sheet_name}'")
            wb.close()
            return "\n".join(lines)[:MAX_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    # ── Стратегия 3: CSV ─────────────────────────────────────────────────────
    if ext == "csv":
        filename = pathlib.Path(path).name
        for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
            try:
                with open(path, "r", encoding=enc, errors="strict", newline="") as f:
                    reader = csv.reader(f)
                    rows = [next(reader, []), next(reader, []), next(reader, [])]
                header = rows[0]
                if not header:
                    return f"CSV файл: {filename}"
                lines = [
                    f"CSV файл: {filename}",
                    f"Колонки: {', '.join(str(v) for v in header if v)}",
                ]
                for row in rows[1:]:
                    if row:
                        lines.append(", ".join(str(v) for v in row if v))
                return "\n".join(lines)[:MAX_CHARS]
            except (UnicodeDecodeError, LookupError):
                continue
            except Exception:
                break
        return f"CSV файл: {filename}"

    # ── Стратегия 3: PPTX ────────────────────────────────────────────────────
    if ext == "pptx":
        try:
            from pptx import Presentation
            prs  = Presentation(path)
            filename = pathlib.Path(path).name
            slide_titles = []
            max_slides = _pages_for_size(path) * 5  # ~5 слайдов на "страницу"
            for slide in prs.slides[:max_slides]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # Первый непустой параграф = вероятный заголовок слайда
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and len(t) > 2:
                                slide_titles.append(t)
                                break
                        break   # берём только первый текстовый блок слайда
            description = f"Презентация: {filename}"
            if slide_titles:
                description += "\nСлайды: " + " | ".join(slide_titles[:15])
            return description[:MAX_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    return ""


# ── OpenAI Embeddings ─────────────────────────────────────────────────────────

def _embed_batch(texts: list[str], api_key: str) -> list[list[float]] | None:
    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key, timeout=30.0, max_retries=1)
        resp   = client.embeddings.create(model=EMBED_MODEL, input=texts)
        return [item.embedding for item in resp.data]
    except Exception as e:
        print(f"    [semantic] Embedding error: {e}")
        return None


# ── Сериализация ──────────────────────────────────────────────────────────────

def _vec_to_blob(vec: list[float]) -> bytes:
    return struct.pack(f"{len(vec)}f", *vec)


def _blob_to_vec(blob: bytes) -> np.ndarray:
    return np.frombuffer(blob, dtype=np.float32).copy()


_EXT_CATEGORY: dict[str, str] = {
    **{e: "document" for e in ("pdf", "docx", "doc", "txt", "md", "rst", "csv", "pptx", "xlsx", "xls", "odt", "rtf")},
    **{e: "code"     for e in FULL_TEXT},
    **{e: "photo"    for e in ("jpg", "jpeg", "png", "gif", "bmp", "webp", "tiff", "svg")},
    **{e: "music"    for e in ("mp3", "flac", "wav", "aac", "ogg", "m4a")},
    **{e: "video"    for e in ("mp4", "avi", "mkv", "mov", "wmv", "webm")},
    **{e: "archive"  for e in ("zip", "rar", "7z", "tar", "gz")},
}


def _ext_to_category(ext: str) -> str:
    return _EXT_CATEGORY.get(ext.lower(), "other")


# ── SemanticIndexer ───────────────────────────────────────────────────────────

class SemanticIndexer:
    def __init__(self):
        DB_PATH.parent.mkdir(parents=True, exist_ok=True)
        self._lock   = threading.Lock()
        self._conn   = sqlite3.connect(str(DB_PATH), check_same_thread=False)
        self._conn.row_factory = sqlite3.Row
        self._init_db()
        self._progress: dict = {
            "is_indexing": False,
            "indexed": 0, "total": 0, "percent": 100,
        }
        # Очередь для фоновой индексации из watchdog
        self._queue: queue.Queue = queue.Queue()
        threading.Thread(
            target=self._background_worker,
            daemon=True,
            name="semantic-worker",
        ).start()

    # ── БД ────────────────────────────────────────────────────────────────────

    def _init_db(self):
        self._conn.executescript("""
            CREATE TABLE IF NOT EXISTS embeddings (
                path         TEXT PRIMARY KEY,
                modified_at  REAL NOT NULL,
                embedding    BLOB NOT NULL,
                text_preview TEXT,
                indexed_at   REAL NOT NULL
            );
            CREATE INDEX IF NOT EXISTS idx_sem_mtime ON embeddings(modified_at);
        """)
        self._conn.commit()

    # ── Фоновый worker (обрабатывает очередь из watchdog) ─────────────────────

    def _background_worker(self):
        """Забирает пути из очереди батчами, ждёт между батчами (rate limit)."""
        while True:
            try:
                path = self._queue.get(timeout=5)
            except queue.Empty:
                continue

            batch = [path]
            # Сливаем всё что успело накопиться
            try:
                while len(batch) < BATCH_SIZE:
                    batch.append(self._queue.get_nowait())
            except queue.Empty:
                pass

            import config
            api_key = getattr(config, "OPENAI_API_KEY", "")
            if api_key:
                self.build_index(batch, api_key)

            time.sleep(1)  # пауза между батчами — OpenAI rate limit

    def enqueue(self, path: str):
        """Добавить файл в очередь на семантическую индексацию (из watchdog)."""
        ext = pathlib.Path(path).suffix.lower().lstrip(".")
        if ext in ALL_SUPPORTED and not _too_large(path):
            self._queue.put(path)

    # ── Построение индекса ────────────────────────────────────────────────────

    def build_index(self, file_paths: list[str], api_key: str) -> int:
        """
        Индексирует переданные файлы.
        Пропускает файлы которые уже в БД и не изменились.
        """
        if not api_key:
            return 0

        # Фильтруем: только поддерживаемые расширения + не слишком большие
        candidates = [
            p for p in file_paths
            if pathlib.Path(p).suffix.lower().lstrip(".") in ALL_SUPPORTED
            and os.path.isfile(p)
            and not _too_large(p)
        ]

        # Только новые или изменённые
        to_index: list[tuple[str, float]] = []
        for path in candidates:
            try:
                mtime = os.path.getmtime(path)
            except OSError:
                continue
            with self._lock:
                row = self._conn.execute(
                    "SELECT modified_at FROM embeddings WHERE path = ?", (path,)
                ).fetchone()
            if row is None or abs(float(row["modified_at"]) - mtime) > 0.5:
                to_index.append((path, mtime))

        if not to_index:
            return 0

        print(f"    [semantic] Индексирую {len(to_index)} файлов...")
        self._progress = {
            "is_indexing": True,
            "indexed": 0, "total": len(to_index), "percent": 0,
        }

        indexed      = 0
        batch_texts:  list[str]   = []
        batch_paths:  list[str]   = []
        batch_mtimes: list[float] = []

        for path, mtime in to_index:
            text = _extract_text(path)
            if not text.strip():
                continue

            # Имя файла в тексте улучшает смысловое совпадение
            stem = pathlib.Path(path).stem.replace("_", " ").replace("-", " ")
            batch_texts.append(f"{stem}. {text}")
            batch_paths.append(path)
            batch_mtimes.append(mtime)

            if len(batch_texts) >= BATCH_SIZE:
                indexed += self._flush_batch(
                    batch_texts, batch_paths, batch_mtimes, api_key
                )
                batch_texts.clear()
                batch_paths.clear()
                batch_mtimes.clear()
                self._progress["indexed"] = indexed
                self._progress["percent"] = int(indexed * 100 / len(to_index))

        if batch_texts:
            indexed += self._flush_batch(
                batch_texts, batch_paths, batch_mtimes, api_key
            )

        self._progress = {
            "is_indexing": False,
            "indexed": indexed, "total": len(to_index), "percent": 100,
        }
        print(f"    [semantic] Готово: {indexed} файлов")
        return indexed

    def _flush_batch(
        self,
        texts:  list[str],
        paths:  list[str],
        mtimes: list[float],
        api_key: str,
    ) -> int:
        vecs = _embed_batch(texts, api_key)
        if not vecs:
            return 0
        now  = time.time()
        rows = [
            (path, mtime, _vec_to_blob(vec), texts[i][:300], now)
            for i, (path, mtime, vec) in enumerate(zip(paths, mtimes, vecs))
        ]
        with self._lock:
            self._conn.executemany(
                "INSERT OR REPLACE INTO embeddings "
                "(path, modified_at, embedding, text_preview, indexed_at) "
                "VALUES (?, ?, ?, ?, ?)",
                rows,
            )
            self._conn.commit()
        return len(rows)

    # ── Поиск ─────────────────────────────────────────────────────────────────

    def search(
        self,
        query:    str,
        api_key:  str,
        limit:    int = 5,
        category: str = "",
    ) -> list[dict]:
        """
        Семантический поиск по содержимому файлов.
        Возвращает список файлов с полем 'score' (0–1) и 'preview'.
        """
        if not api_key or not query.strip():
            return []

        vecs = _embed_batch([query], api_key)
        if not vecs:
            return []
        query_vec = np.array(vecs[0], dtype=np.float32)

        with self._lock:
            rows = self._conn.execute(
                "SELECT path, embedding, text_preview FROM embeddings"
            ).fetchall()

        if not rows:
            return []

        live_paths: list[str]        = []
        matrix:     list[np.ndarray] = []
        previews:   dict[str, str]   = {}

        for r in rows:
            p = r["path"]
            if not os.path.isfile(p):   # исключаем папки и несуществующие пути
                continue
            ext = pathlib.Path(p).suffix.lower().lstrip(".")
            # Фильтр по категории
            if category:
                if category == "document" and ext not in (
                    set(LIBRARY) | {"txt", "md", "rst", "csv", "pptx", "xlsx"}
                ):
                    continue
                elif category == "code" and ext not in FULL_TEXT:
                    continue
                elif category == "photo" and ext not in {
                    "jpg", "jpeg", "png", "gif", "bmp", "webp", "tiff", "svg"
                }:
                    continue
                elif category == "music" and ext not in {
                    "mp3", "flac", "wav", "aac", "ogg", "m4a"
                }:
                    continue
                elif category == "video" and ext not in {
                    "mp4", "avi", "mkv", "mov", "wmv", "webm"
                }:
                    continue
            live_paths.append(p)
            matrix.append(_blob_to_vec(r["embedding"]))
            previews[p] = r["text_preview"] or ""

        if not matrix:
            return []

        mat       = np.array(matrix, dtype=np.float32)
        mat_norms = np.linalg.norm(mat, axis=1)
        q_norm    = float(np.linalg.norm(query_vec))
        if q_norm < 1e-9:
            return []

        sims    = (mat @ query_vec) / (mat_norms * q_norm + 1e-9)
        top_idx = np.argsort(sims)[::-1]

        import datetime
        results = []
        for idx in top_idx:
            score = float(sims[idx])
            if score < SIM_THRESH:
                break
            path = live_paths[idx]
            p    = pathlib.Path(path)
            try:
                mtime     = os.path.getmtime(path)
                mtime_str = datetime.datetime.fromtimestamp(mtime).strftime("%d.%m.%Y %H:%M")
                size      = os.path.getsize(path)
            except OSError:
                mtime_str = ""
                size      = 0

            file_ext = p.suffix.lower().lstrip(".")
            results.append({
                "id":             None,
                "name":           p.name,
                "path":           path,
                "folder":         str(p.parent),
                "extension":      file_ext,
                "category":       _ext_to_category(file_ext),
                "size_bytes":     size,
                "size_human":     _human_size(size),
                "modified_at":    0,
                "modified_human": mtime_str,
                "score":          round(score, 3),
                "preview":        previews[path][:200],
            })
            if len(results) >= limit:
                break

        return results

    # ── Управление ────────────────────────────────────────────────────────────

    def remove_file(self, path: str):
        with self._lock:
            self._conn.execute(
                "DELETE FROM embeddings WHERE path = ?", (path,)
            )
            self._conn.commit()

    def get_status(self) -> dict:
        with self._lock:
            count = self._conn.execute(
                "SELECT COUNT(*) FROM embeddings"
            ).fetchone()[0]
        return {"indexed_files": count, **self._progress}


def _human_size(b: int) -> str:
    if b < 1024:       return f"{b} Б"
    if b < 1024 ** 2:  return f"{b // 1024} КБ"
    if b < 1024 ** 3:  return f"{b // 1024 ** 2} МБ"
    return f"{b / 1024 ** 3:.1f} ГБ"


# ── Синглтон ──────────────────────────────────────────────────────────────────

_instance: "SemanticIndexer | None" = None
_inst_lock = threading.Lock()


def get_semantic_indexer() -> SemanticIndexer:
    global _instance
    if _instance is None:
        with _inst_lock:
            if _instance is None:
                _instance = SemanticIndexer()
    return _instance
